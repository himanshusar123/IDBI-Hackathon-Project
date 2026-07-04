import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    # Widescreen 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme colors matching the app aesthetics
    COLOR_BG = RGBColor(8, 12, 20)          # Dark Slate (#080c14)
    COLOR_CYAN = RGBColor(6, 182, 212)      # Electric Cyan (#06b6d4)
    COLOR_PURPLE = RGBColor(139, 92, 246)   # Soft Purple (#8b5cf6)
    COLOR_TEXT = RGBColor(248, 250, 252)    # Near White (#f8fafc)
    COLOR_MUTED = RGBColor(148, 163, 184)   # Gray (#94a3b8)
    COLOR_CARD = RGBColor(22, 30, 49)       # Card BG (#161e31)

    # 1. Slide creation helper
    def add_custom_slide(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
        
        # Set solid dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG
        
        # Add slide title (if provided)
        if title_text:
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1.0))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = 'Outfit'
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN
            
            # Subtle title bottom line
            line = slide.shapes.add_shape(
                1, # rectangle
                Inches(0.75), Inches(1.4), Inches(2.0), Inches(0.04)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = COLOR_PURPLE
            line.line.fill.background()
            
        return slide

    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG
    
    # Title & Subtitle in single text frame to avoid overlap
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = "AURA WEALTH"
    p_title.font.name = 'Outfit'
    p_title.font.size = Pt(64)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_CYAN
    p_title.alignment = PP_ALIGN.LEFT
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Next-Gen AI-Powered Avatar Wealth Advisor for Modern Banking"
    p_sub.font.name = 'Inter'
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = COLOR_PURPLE
    p_sub.space_before = Pt(12)
    
    p_meta = tf.add_paragraph()
    p_meta.text = "Personalized • Scalable • Conversational Wealth Advisory"
    p_meta.font.name = 'Inter'
    p_meta.font.size = Pt(14)
    p_meta.font.color.rgb = COLOR_MUTED
    p_meta.space_before = Pt(30)

    # ==================== SLIDE 2: PROBLEM STATEMENT ====================
    slide2 = add_custom_slide("The Fragmented Wealth Management Problem")
    
    # Left Column (Challenge description)
    left_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "Key Market Pain Points:"
    p.font.name = 'Outfit'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    
    bullets = [
      "Advisory Is Inaccessible: High net worth advisory services are expensive and locked behind high capital minimums, leaving retail customers underserved.",
      "Fragmented Data: Investment behaviors, daily transactions, and savings habits are scattered, restricting banks from providing comprehensive advice.",
      "Lack of Timely Guidance: Manual reviews are slow. Users need instantaneous, proactive alerts on budget drifts and portfolio rebalancing."
    ]
    for b in bullets:
        p_b = tf_left.add_paragraph()
        p_b.text = "• " + b.split(":")[0] + ":" + b.split(":")[1]
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = COLOR_MUTED
        p_b.space_before = Pt(14)

    # Right Column (Interactive Panel Visual Mockup)
    right_card = slide2.shapes.add_shape(1, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_CARD
    right_card.line.color.rgb = COLOR_PURPLE
    right_card.line.width = Pt(1.5)
    
    tf_card = right_card.text_frame
    tf_card.word_wrap = True
    tf_card.margin_left = Inches(0.4)
    tf_card.margin_right = Inches(0.4)
    tf_card.margin_top = Inches(0.4)
    
    p_card = tf_card.paragraphs[0]
    p_card.text = "THE IMPACT OF INACTION"
    p_card.font.name = 'Outfit'
    p_card.font.size = Pt(18)
    p_card.font.bold = True
    p_card.font.color.rgb = COLOR_CYAN
    
    p_card_body = tf_card.add_paragraph()
    p_card_body.text = "Traditional banks lose customer engagement to third-party robo-advisors. Without unified wealth intelligence, banks cannot effectively upsell wealth products (mutual funds, deposits) because they lack dynamic, context-aware user profiles."
    p_card_body.font.name = 'Inter'
    p_card_body.font.size = Pt(15)
    p_card_body.font.color.rgb = COLOR_TEXT
    p_card_body.space_before = Pt(14)

    # ==================== SLIDE 3: SOLUTION OVERVIEW ====================
    slide3 = add_custom_slide("Introducing Aura: The AI Advisor")
    
    # 3 pillars
    pillars = [
      {"title": "Conversational Avatar Interface", "desc": "Interactive 3D/SVG visual advisor ('Aura') featuring lip-sync and auditory feedback, providing human-like advisory scaled to millions of banking customers.", "icon": "Avatar"},
      {"title": "Unified Wealth Analytics", "desc": "Integrates transaction history, savings rate, and credit debt to model holistic spending insights and dynamic wealth trajectories.", "icon": "Analytics"},
      {"title": "Automated Asset Rebalancing", "desc": "One-click drift correction that realigns portfolios back to target models dynamically computed from active risk profile assessments.", "icon": "Execution"}
    ]
    
    for i, pil in enumerate(pillars):
        x = Inches(0.75 + i * 4.0)
        card = slide3.shapes.add_shape(1, x, Inches(2.0), Inches(3.7), Inches(4.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_CYAN if i==0 else COLOR_MUTED
        card.line.width = Pt(1)
        
        tf_pil = card.text_frame
        tf_pil.word_wrap = True
        tf_pil.margin_left = Inches(0.3)
        tf_pil.margin_right = Inches(0.3)
        tf_pil.margin_top = Inches(0.3)
        
        p_t = tf_pil.paragraphs[0]
        p_t.text = f"0{i+1}. {pil['title']}"
        p_t.font.name = 'Outfit'
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CYAN
        
        p_d = tf_pil.add_paragraph()
        p_d.text = pil['desc']
        p_d.font.name = 'Inter'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = COLOR_TEXT
        p_d.space_before = Pt(14)

    # ==================== SLIDE 4: SYSTEM ARCHITECTURE ====================
    slide4 = add_custom_slide("Technical Architecture")
    
    # Left Box - Frontend UI
    fe_card = slide4.shapes.add_shape(1, Inches(0.75), Inches(1.8), Inches(3.6), Inches(4.5))
    fe_card.fill.solid()
    fe_card.fill.fore_color.rgb = COLOR_CARD
    fe_card.line.color.rgb = COLOR_MUTED
    tf_fe = fe_card.text_frame
    tf_fe.word_wrap = True
    tf_fe.margin_left = Inches(0.3)
    p_fe = tf_fe.paragraphs[0]
    p_fe.text = "FRONTEND (MOBILE INTERACTION)"
    p_fe.font.name = 'Outfit'
    p_fe.font.size = Pt(16)
    p_fe.font.bold = True
    p_fe.font.color.rgb = COLOR_CYAN
    fe_bullets = [
      "Vite + React SPA structure",
      "High-performance custom SVG animations for lipsync/state indicators",
      "Web Speech API integration for local Speech-to-Text & Text-to-Speech",
      "Custom responsive CSS dashboard"
    ]
    for b in fe_bullets:
        p_b = tf_fe.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_MUTED
        p_b.space_before = Pt(8)

    # Middle Box - AI Advisory Brain
    ai_card = slide4.shapes.add_shape(1, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.5))
    ai_card.fill.solid()
    ai_card.fill.fore_color.rgb = COLOR_CARD
    ai_card.line.color.rgb = COLOR_PURPLE
    ai_card.line.width = Pt(1.5)
    tf_ai = ai_card.text_frame
    tf_ai.word_wrap = True
    tf_ai.margin_left = Inches(0.3)
    p_ai = tf_ai.paragraphs[0]
    p_ai.text = "AI ENGINE (GEMINI PRO / VERTEX)"
    p_ai.font.name = 'Outfit'
    p_ai.font.size = Pt(16)
    p_ai.font.bold = True
    p_ai.font.color.rgb = COLOR_PURPLE
    ai_bullets = [
      "Gemini Flash 2.5 API integration",
      "System context framing models: certified personal banker profile",
      "Real-time data injection (spending history, risk metrics)",
      "Structured JSON responses for automated goal calculations"
    ]
    for b in ai_bullets:
        p_b = tf_ai.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_MUTED
        p_b.space_before = Pt(8)

    # Right Box - Core Bank Integration
    bk_card = slide4.shapes.add_shape(1, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.5))
    bk_card.fill.solid()
    bk_card.fill.fore_color.rgb = COLOR_CARD
    bk_card.line.color.rgb = COLOR_MUTED
    tf_bk = bk_card.text_frame
    tf_bk.word_wrap = True
    tf_bk.margin_left = Inches(0.3)
    p_bk = tf_bk.paragraphs[0]
    p_bk.text = "CORE BANK INTEGRATION"
    p_bk.font.name = 'Outfit'
    p_bk.font.size = Pt(16)
    p_bk.font.bold = True
    p_bk.font.color.rgb = COLOR_CYAN
    bk_bullets = [
      "Open Banking APIs for live balance syncing",
      "Risk profiling engine questionnaire triggers",
      "Brokerage API calls for rebalancing trade executions",
      "Notification triggers for budget deviations"
    ]
    for b in bk_bullets:
        p_b = tf_bk.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_MUTED
        p_b.space_before = Pt(8)

    # ==================== SLIDE 5: VALUE FOR BANKS ====================
    slide5 = add_custom_slide("Business & Value Proposition")
    
    # 2 columns layout
    l_box = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_l = l_box.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "Value for Banking Customers:"
    p_l.font.name = 'Outfit'
    p_l.font.size = Pt(20)
    p_l.font.bold = True
    p_l.font.color.rgb = COLOR_TEXT
    
    c_benefits = [
      "Hyper-Personalization: Guidance calibrated directly to their personal transactions rather than generic articles.",
      "Increased Accessibility: Demystifies investment terms via natural chat or voice discussion.",
      "Clear Actionable Steps: Simplifies rebalancing and budget adjustments into single-tap workflows."
    ]
    for b in c_benefits:
        p_b = tf_l.add_paragraph()
        p_b.text = "• " + b.split(":")[0] + ":" + b.split(":")[1]
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = COLOR_MUTED
        p_b.space_before = Pt(12)

    r_box = slide5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    tf_r = r_box.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "Value for the Financial Institution:"
    p_r.font.name = 'Outfit'
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = COLOR_TEXT
    
    b_benefits = [
      "Drive Asset growth (AUM): Seamlessly guide excess cash into appropriate mutual funds and fixed deposits.",
      "Boost Retention: Interactive features and helpful notifications foster persistent daily banking interactions.",
      "Reduce Operational Costs: Minimizes dependencies on human relationship managers for routine portfolio checkups."
    ]
    for b in b_benefits:
        p_b = tf_r.add_paragraph()
        p_b.text = "• " + b.split(":")[0] + ":" + b.split(":")[1]
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = COLOR_MUTED
        p_b.space_before = Pt(12)

    # Save presentation
    filename = "C:\\Users\\Himanshu Sardana\\.gemini\\antigravity-ide\\scratch\\digital-wealth-avatar\\Aura_Wealth_Advisor_Presentation.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully to {filename}")

if __name__ == "__main__":
    create_presentation()
